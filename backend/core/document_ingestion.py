from __future__ import annotations

import csv
import importlib.util
import json
import mimetypes
import shutil
from dataclasses import asdict, dataclass, field
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, Iterable

from backend.core import observability, safeguards

MAX_FILE_BYTES = int(25 * 1024 * 1024)
MAX_EXTRACTED_CHARS = int(500_000)
logger = observability.get_logger(__name__)


@dataclass
class DocumentPart:
    text: str
    kind: str = "text"
    page: int | None = None
    sheet: str | None = None
    slide: int | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentExtraction:
    filename: str
    mime_type: str
    text: str
    parts: list[DocumentPart]
    success: bool
    warnings: list[str] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_dict(self, include_parts: bool = True) -> dict[str, Any]:
        payload = asdict(self)
        if not include_parts:
            payload.pop("parts", None)
        payload["chars"] = len(self.text)
        payload["words"] = len(self.text.split())
        payload["preview"] = self.text[:3000]
        return payload

    def as_prompt_block(self) -> str:
        header = f"--- FILE: {self.filename} ---"
        meta = [
            f"type={self.mime_type or 'unknown'}",
            f"chars={len(self.text)}",
            f"words={len(self.text.split())}",
        ]
        if self.warnings:
            meta.append("warnings=" + "; ".join(self.warnings[:3]))
        body = self.text or "[No text could be extracted from this file.]"
        return "\n" + header + "\n" + "[" + " | ".join(meta) + "]\n" + body


def capabilities() -> dict[str, Any]:
    packages = {
        "pdfplumber": _has_package("pdfplumber"),
        "pymupdf": _has_package("fitz"),
        "python_docx": _has_package("docx"),
        "openpyxl": _has_package("openpyxl"),
        "python_pptx": _has_package("pptx"),
        "pillow": _has_package("PIL"),
        "pytesseract": _has_package("pytesseract"),
        "bs4": _has_package("bs4"),
    }
    binaries = {
        "tesseract": bool(shutil.which("tesseract")),
        "ocrmypdf": bool(shutil.which("ocrmypdf")),
    }
    return {
        "max_file_mb": round(MAX_FILE_BYTES / 1024 / 1024, 1),
        "max_extracted_chars": MAX_EXTRACTED_CHARS,
        "packages": packages,
        "binaries": binaries,
        "supported_extensions": [
            ".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".txt", ".md",
            ".json", ".ipynb", ".html", ".htm", ".png", ".jpg", ".jpeg", ".webp",
        ],
        "ocr_ready": packages["pytesseract"] and packages["pillow"] and binaries["tesseract"],
        "layout_pdf_ready": packages["pdfplumber"] or packages["pymupdf"],
        "office_ready": packages["python_docx"] and packages["openpyxl"],
    }


def extract_uploads(files: Iterable[Any]) -> list[DocumentExtraction]:
    results = []
    for upload in files:
        filename = getattr(upload, "filename", None) or "uploaded_file"
        raw = upload.file.read()
        try:
            upload.file.seek(0)
        except Exception:
            pass
        results.append(extract_bytes(filename, raw, getattr(upload, "content_type", None)))
    return results


def extract_path(path: str | Path) -> DocumentExtraction:
    file_path = Path(path)
    return extract_bytes(file_path.name, file_path.read_bytes(), mimetypes.guess_type(file_path.name)[0])


def extract_bytes(filename: str, raw: bytes, content_type: str | None = None) -> DocumentExtraction:
    warnings: list[str] = []
    if len(raw) > MAX_FILE_BYTES:
        raw = raw[:MAX_FILE_BYTES]
        warnings.append("File was truncated before extraction because it exceeded the local size limit.")

    ext = Path(filename).suffix.lower()
    mime_type = content_type or mimetypes.guess_type(filename)[0] or ""
    parts: list[DocumentPart] = []

    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            parts, warnings = _extract_pdf(raw, warnings)
        elif ext == ".docx":
            parts, warnings = _extract_docx(raw, warnings)
        elif ext == ".xlsx":
            parts, warnings = _extract_xlsx(raw, warnings)
        elif ext == ".pptx":
            parts, warnings = _extract_pptx(raw, warnings)
        elif ext in {".png", ".jpg", ".jpeg", ".webp"} or mime_type.startswith("image/"):
            parts, warnings = _extract_image(raw, warnings)
        elif ext in {".html", ".htm"} or mime_type == "text/html":
            parts = [DocumentPart(_html_to_text(_decode_text(raw)), kind="html")]
        elif ext == ".json" or mime_type == "application/json":
            parts = [DocumentPart(_json_to_text(raw, warnings), kind="json")]
        elif ext == ".ipynb":
            parts = _extract_ipynb(raw, warnings)
        elif ext == ".csv" or mime_type == "text/csv":
            parts = [DocumentPart(_csv_to_text(raw, warnings), kind="table")]
        else:
            parts = [DocumentPart(_decode_text(raw), kind="text")]
    except Exception as exc:
        warnings.append(f"Extraction failed: {exc}")
        parts = []

    text = _join_parts(parts)
    text = safeguards.truncate_text(text, MAX_EXTRACTED_CHARS, "extracted document")
    success = bool(text.strip())
    if not success and not warnings:
        warnings.append("No readable text was found. The file may be scanned, encrypted, or image-only.")

    observability.log_event(
        logger,
        "document.extract",
        filename=filename,
        mime_type=mime_type,
        success=success,
        chars=len(text),
        warnings=warnings[:3],
    )
    return DocumentExtraction(
        filename=filename,
        mime_type=mime_type,
        text=text,
        parts=parts,
        success=success,
        warnings=warnings,
        metadata={"extension": ext, "bytes": len(raw)},
    )


def _extract_pdf(raw: bytes, warnings: list[str]) -> tuple[list[DocumentPart], list[str]]:
    parts: list[DocumentPart] = []
    if _has_package("pdfplumber"):
        try:
            import pdfplumber

            with pdfplumber.open(BytesIO(raw)) as pdf:
                for page_index, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        parts.append(DocumentPart(page_text, kind="page", page=page_index))
                    try:
                        for table_index, table in enumerate(page.extract_tables() or [], start=1):
                            table_text = _table_to_markdown(table)
                            if table_text:
                                parts.append(
                                    DocumentPart(
                                        table_text,
                                        kind="table",
                                        page=page_index,
                                        metadata={"table": table_index},
                                    )
                                )
                    except Exception as exc:
                        warnings.append(f"PDF table extraction failed on page {page_index}: {exc}")
        except Exception as exc:
            warnings.append(f"pdfplumber failed: {exc}")

    if not parts and _has_package("fitz"):
        try:
            import fitz

            doc = fitz.open(stream=raw, filetype="pdf")
            for page_index, page in enumerate(doc, start=1):
                page_text = page.get_text("text") or ""
                if page_text.strip():
                    parts.append(DocumentPart(page_text, kind="page", page=page_index))
        except Exception as exc:
            warnings.append(f"PyMuPDF failed: {exc}")

    if not parts:
        ocr_parts, warnings = _ocr_pdf(raw, warnings)
        parts.extend(ocr_parts)
        if not ocr_parts:
            warnings.append("No PDF text was extracted. Install OCR dependencies for scanned PDFs.")
    return parts, warnings


def _extract_docx(raw: bytes, warnings: list[str]) -> tuple[list[DocumentPart], list[str]]:
    if not _has_package("docx"):
        return [], warnings + ["Missing python-docx; cannot extract DOCX files."]
    import docx

    doc = docx.Document(BytesIO(raw))
    parts: list[DocumentPart] = []
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if paragraphs:
        parts.append(DocumentPart("\n".join(paragraphs), kind="text"))
    for table_index, table in enumerate(doc.tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        table_text = _table_to_markdown(rows)
        if table_text:
            parts.append(DocumentPart(table_text, kind="table", metadata={"table": table_index}))
    return parts, warnings


def _extract_xlsx(raw: bytes, warnings: list[str]) -> tuple[list[DocumentPart], list[str]]:
    if not _has_package("openpyxl"):
        return [], warnings + ["Missing openpyxl; cannot extract XLSX files."]
    import openpyxl

    workbook = openpyxl.load_workbook(BytesIO(raw), data_only=False, read_only=True)
    parts: list[DocumentPart] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                rows.append(values)
            if len(rows) >= 500:
                warnings.append(f"Sheet '{sheet.title}' was truncated at 500 non-empty rows.")
                break
        table_text = _table_to_markdown(rows)
        if table_text:
            parts.append(DocumentPart(table_text, kind="sheet", sheet=sheet.title))
    return parts, warnings


def _extract_ipynb(raw: bytes, warnings: list[str]) -> list[DocumentPart]:
    try:
        notebook = json.loads(_decode_text(raw))
    except Exception as exc:
        warnings.append(f"Notebook JSON parsing failed: {exc}")
        return []

    parts: list[DocumentPart] = []
    cells = notebook.get("cells", [])
    if not isinstance(cells, list):
        warnings.append("Notebook did not contain a valid cells list.")
        return []

    for index, cell in enumerate(cells, start=1):
        if not isinstance(cell, dict):
            continue
        cell_type = str(cell.get("cell_type") or "cell")
        source = cell.get("source", "")
        if isinstance(source, list):
            source_text = "".join(str(line) for line in source)
        else:
            source_text = str(source or "")
        source_text = source_text.strip()
        if not source_text:
            continue
        header = f"Cell {index} [{cell_type}]"
        if cell_type == "code":
            text = header + "\n```python\n" + source_text + "\n```"
        else:
            text = header + "\n" + source_text
        parts.append(DocumentPart(text, kind="notebook_cell", metadata={"cell": index, "cell_type": cell_type}))

    return parts


def _extract_pptx(raw: bytes, warnings: list[str]) -> tuple[list[DocumentPart], list[str]]:
    if not _has_package("pptx"):
        return [], warnings + ["Missing python-pptx; cannot extract PPTX files."]
    from pptx import Presentation

    presentation = Presentation(BytesIO(raw))
    parts: list[DocumentPart] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                table_text = _table_to_markdown(rows)
                if table_text:
                    texts.append(table_text)
        if texts:
            parts.append(DocumentPart("\n".join(texts), kind="slide", slide=slide_index))
    return parts, warnings


def _extract_image(raw: bytes, warnings: list[str]) -> tuple[list[DocumentPart], list[str]]:
    if not (_has_package("PIL") and _has_package("pytesseract") and shutil.which("tesseract")):
        return [], warnings + ["OCR is not ready; install Pillow, pytesseract, and the tesseract binary to read images."]
    from PIL import Image
    import pytesseract

    text = pytesseract.image_to_string(Image.open(BytesIO(raw))).strip()
    return ([DocumentPart(text, kind="ocr_image")] if text else []), warnings


def _ocr_pdf(raw: bytes, warnings: list[str]) -> tuple[list[DocumentPart], list[str]]:
    if not (_has_package("fitz") and _has_package("PIL") and _has_package("pytesseract") and shutil.which("tesseract")):
        return [], warnings
    import fitz
    import pytesseract
    from PIL import Image

    parts: list[DocumentPart] = []
    doc = fitz.open(stream=raw, filetype="pdf")
    max_pages = min(len(doc), int(10))
    for page_index in range(max_pages):
        page = doc[page_index]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.open(BytesIO(pix.tobytes("png")))
        text = pytesseract.image_to_string(image).strip()
        if text:
            parts.append(DocumentPart(text, kind="ocr_page", page=page_index + 1))
    if len(doc) > max_pages:
        warnings.append(f"OCR was limited to the first {max_pages} pages.")
    return parts, warnings


def _csv_to_text(raw: bytes, warnings: list[str]) -> str:
    decoded = _decode_text(raw)
    rows = list(csv.reader(StringIO(decoded)))
    if len(rows) > 500:
        rows = rows[:500]
        warnings.append("CSV was truncated at 500 rows.")
    return _table_to_markdown(rows)


def _json_to_text(raw: bytes, warnings: list[str]) -> str:
    try:
        parsed = json.loads(_decode_text(raw))
        return json.dumps(parsed, indent=2, ensure_ascii=False)
    except Exception as exc:
        warnings.append(f"JSON parse failed, using raw text: {exc}")
        return _decode_text(raw)


def _html_to_text(html: str) -> str:
    if _has_package("bs4"):
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return " ".join(soup.get_text(" ").split())
    return " ".join(html.split())


def _table_to_markdown(rows: list[list[Any]]) -> str:
    cleaned = [["" if value is None else str(value).replace("\n", " ").strip() for value in row] for row in rows]
    cleaned = [row for row in cleaned if any(cell for cell in row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    padded = [row + [""] * (width - len(row)) for row in cleaned]
    header = padded[0]
    body = padded[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _join_parts(parts: list[DocumentPart]) -> str:
    blocks = []
    for part in parts:
        label = part.kind.upper()
        if part.page is not None:
            label += f" page={part.page}"
        if part.sheet:
            label += f" sheet={part.sheet}"
        if part.slide is not None:
            label += f" slide={part.slide}"
        blocks.append(f"[{label}]\n{part.text.strip()}")
    return "\n\n".join(block for block in blocks if block.strip())


def _decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _has_package(name: str) -> bool:
    return importlib.util.find_spec(name) is not None
